from pptx.parts.image import Image, ImagePart
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.oxml.shapes.picture import CT_Picture
from pptx.shapes.placeholder import PlaceholderPicture
from ..style import getAllStylingAreasPosInParagraphs, removeStylingAreas, parser, basicFormating, getStylingAreasPosInTextFrame
import re
import base64

def isImageBased(placeholder):
  """
  Returns true is the placeholder is "image-based"

  Args:
    placeholder: PicturePlaceholder
  """
  return placeholder.placeholder_format.type._member_name in ('PICTURE')

def imagePartFromBlob(placeholder, blob):
  """
  Returns the image part corresponding to the image

  Args:
    placeholder: Picture Placeholder
    blob (bytes): blob image
  """

  # Check if the image doesn't already exists
  parts = placeholder.part.package._image_parts
  image = Image.from_blob(blob)

  imagePart = parts._find_by_sha1(image.sha1)

  if imagePart is None:
    # If not create a new ImagePart
    imagePart = ImagePart.new(parts._package, image)

  imagePart = ImagePart.new(parts._package, image)

  return imagePart

def imagePartAndrIdFromBlob(placeholder, blob):
  """
  Returns the imagePart and the rId corresponding to the image

  Args:
    placeholder: Picture Placeholder
    blob (bytes): blob image
  """
  imagePart = imagePartFromBlob(placeholder, blob)
  rId = placeholder.part.relate_to(imagePart, RT.IMAGE)

  return imagePart, rId

def insertBlobImage(placeholder, blob, objectFit):
  """
  Inserts a blob image into a placeholder

  Args:
    placeholder: Picture Placeholder
    blob (bytes): blob image
  """
  # TODO: Multiple possibilities for resizing
  imagePart, rId = imagePartAndrIdFromBlob(placeholder, blob)
  desc, imageSize = imagePart.desc, imagePart._px_size

  shapeId, name = placeholder.shape_id, placeholder.name

  pic = CT_Picture.new_ph_pic(shapeId, name, desc, rId)

  if(objectFit == "contain"):
    pic.blipFill.crop(_contain_cropping(imageSize, (placeholder.width, placeholder.height)))
  else:
    pic.crop_to_fit(imageSize, (placeholder.width, placeholder.height))

  placeholder._replace_placeholder_with(pic)

  return PlaceholderPicture(pic, placeholder._parent)

def formateImagePlaceholder(placeholder, data):
  """
  Place an image in the placeholder according to the first tag found

  Args:
    placeholder: Picture Placeholder
    blob (bytes): blob image
  
  Returns the PlaceholderPicture object of the image
  """
  # Parse tex_frame
  texts = [paragraph.text for paragraph in placeholder._base_placeholder.text_frame.paragraphs]
  stylingAreas = list(getStylingAreasPosInTextFrame(texts))
  style = parser(stylingAreas, data)
  
  text = ''.join(texts)
  # Remove styling areas
  for stylingArea in stylingAreas:
    text = text[:stylingArea.start()] + text[stylingArea.end():]

  tag = re.search(r"\$\{([A-Za-z0-9._\-]+)\}", text)

  if tag:
    tag = tag.group(1)

  if(tag and data.get(tag)):
    blob = base64.decodestring(bytes(data.get(tag), 'utf-8'))
    placeholderPicture = insertBlobImage(placeholder, blob, style.get('object-fit'))

    # Apply style at the end because placeholder is replaced
    basicFormating(placeholderPicture, style)
    return placeholderPicture

  basicFormating(placeholder, style)

  return placeholder

  

def _contain_cropping(image_size, view_size):
    """
    Return a (left, top, right, bottom) 4-tuple containing the cropping
    values required to display an image of *image_size* in *view_size*
    when stretched proportionately. Each value is a percentage expressed
    as a fraction of 1.0, e.g. 0.425 represents 42.5%. *image_size* and
    *view_size* are each (width, height) pairs.
    """

    def aspect_ratio(width, height):
        return width / height

    ar_view = aspect_ratio(*view_size)
    ar_image = aspect_ratio(*image_size)

    if ar_view > ar_image:  
        crop = (1.0 - (ar_view / ar_image)) / 2.0
        return (crop, 0.0, crop, 0.0)
    if ar_view < ar_image:  
        crop = (1.0 - (ar_image / ar_view)) / 2.0
        return (0.0, crop, 0.0, crop)
    return (0.0, 0.0, 0.0, 0.0)

