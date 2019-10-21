import json
from pptx import Presentation
from pptx_creator.slide import generateSlide

def fillPresentation(presentation, slides):
  """
  Fill a whole presentation

  Args:
    presentation
    slides: a dictionnary containing the configs for all slides
  
  Returns the presentation object
  """

  for slide in slides:
    generateSlide(presentation, slide)
  
  return presentation

def generatePresentation(templatePath, configPath, outPath):
  """
  Creates the presentation and fills it with the config

  Args:
    templatePath (string): the path to the pptx template file
    configPath (string): the path to the json object containing the configuration for all the slides
    outPath (string): the path to the output

  Returns the presentation object
  """

  presentation = Presentation('./personnal/templates/Template4.pptx')

  with open(configPath, encoding='utf-8') as json_file:
    slides = json.load(json_file)

  fillPresentation(presentation, slides)

  presentation.save(outPath)

  return presentation
  